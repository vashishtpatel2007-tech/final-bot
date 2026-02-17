"""
RAG module — ChromaDB vector store with multi-format document ingestion.
Handles: PDF, DOCX, PPTX, XLSX, images (JPEG/PNG), and plain text.
Uses sentence-transformers for local embeddings (free, no API needed).
"""
import os
import tempfile
import chromadb
from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import google.generativeai as genai

# ── ChromaDB Setup ────────────────────────────────────────────

CHROMA_DIR = os.getenv("CHROMA_DIR", "./chroma_db")
embedding_fn = SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")
chroma_client = chromadb.PersistentClient(path=CHROMA_DIR)
collection = chroma_client.get_or_create_collection(
    name="academic_docs",
    embedding_function=embedding_fn,
    metadata={"hnsw:space": "cosine"},
)

# Track ingested files to avoid duplicates
MANIFEST_PATH = os.path.join(CHROMA_DIR, "ingested_files.txt")


def _get_ingested_files() -> set:
    if os.path.exists(MANIFEST_PATH):
        with open(MANIFEST_PATH, "r") as f:
            return set(line.strip() for line in f if line.strip())
    return set()


def _mark_ingested(file_id: str):
    with open(MANIFEST_PATH, "a") as f:
        f.write(file_id + "\n")


# ── File Parsers ──────────────────────────────────────────────

def parse_pdf(file_path: str) -> str:
    """Extract text from PDF. Falls back to Gemini Vision for scanned PDFs."""
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"

    # If very little text extracted, it's likely a scanned PDF → use Gemini Vision
    if len(text.strip()) < 50:
        text = parse_image_with_gemini(file_path, is_pdf=True)

    return text


def parse_docx(file_path: str) -> str:
    """Extract text from DOCX files, including tables."""
    doc = DocxDocument(file_path)
    parts = []

    # Extract paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    # Extract tables
    for table in doc.tables:
        table_text = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_text.append(" | ".join(row_data))
        parts.append("\n".join(table_text))

    return "\n".join(parts)


def parse_image_with_gemini(file_path: str, is_pdf: bool = False) -> str:
    """Use Gemini Vision to extract text/tables from images and scanned PDFs.
    This is the key feature — Gemini can read timetables, handwritten notes, etc."""
    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
    model = genai.GenerativeModel("gemini-2.0-flash")

    prompt = """Extract ALL text content from this image/document. 
    If there are tables, timetables, or structured data, convert them into a clear text format.
    If there are rows and columns, preserve the structure using | separators.
    If there is handwritten text, do your best to transcribe it.
    Include every piece of information visible in the image.
    Do NOT add any commentary — just extract the raw content."""

    if is_pdf:
        # For scanned PDFs, upload the file
        uploaded = genai.upload_file(file_path)
        response = model.generate_content([prompt, uploaded])
    else:
        # For images, open with PIL
        img = Image.open(file_path)
        response = model.generate_content([prompt, img])

    return response.text if response.text else ""


def parse_pptx(file_path: str) -> str:
    """Extract text from PowerPoint PPTX files."""
    prs = Presentation(file_path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    slide_text.append(" | ".join(row_data))
        parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


def parse_xlsx(file_path: str) -> str:
    """Extract text from Excel XLSX files."""
    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_text = [f"--- Sheet: {sheet_name} ---"]
        for row in ws.iter_rows(values_only=True):
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(cell.strip() for cell in row_data):
                sheet_text.append(" | ".join(row_data))
        parts.append("\n".join(sheet_text))
    wb.close()
    return "\n\n".join(parts)


def parse_file(file_path: str) -> str:
    """Auto-detect file type and extract text content."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return parse_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        if ext == ".pptx":
            return parse_pptx(file_path)
        else:
            # .ppt (old format) — try Gemini Vision as fallback
            return parse_image_with_gemini(file_path, is_pdf=True)
    elif ext in (".xlsx", ".xls"):
        if ext == ".xlsx":
            return parse_xlsx(file_path)
        else:
            # .xls (old format) — try as text
            return parse_image_with_gemini(file_path, is_pdf=True)
    elif ext in (".jpg", ".jpeg", ".png", ".bmp", ".gif", ".webp"):
        return parse_image_with_gemini(file_path)
    elif ext in (".txt", ".md", ".csv"):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    else:
        # Try as text file
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""


# ── Chunking ──────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """Split text into overlapping chunks for better retrieval."""
    if not text or len(text.strip()) < 10:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap

    return chunks


# ── Ingest & Query ────────────────────────────────────────────

def ingest_document(file_path: str, stream: str, year: int, doc_type: str = "notes",
                    filename: str = "", drive_link: str = "", file_id: str = ""):
    """Parse a file, chunk it, embed it, and store in ChromaDB with metadata."""

    # Skip if already ingested
    if file_id and file_id in _get_ingested_files():
        print(f"  ⏭️  Already ingested: {filename}")
        return 0

    text = parse_file(file_path)
    if not text.strip():
        print(f"  ⚠️  No text extracted from: {filename}")
        return 0

    chunks = chunk_text(text)
    if not chunks:
        return 0

    # Create unique IDs and metadata for each chunk
    ids = [f"{stream}_{year}_{file_id or filename}_{i}" for i in range(len(chunks))]
    metadatas = [
        {
            "stream": stream.upper(),
            "year": year,
            "type": doc_type,
            "filename": filename,
            "drive_link": drive_link,
        }
        for _ in chunks
    ]

    # Upsert into ChromaDB
    collection.upsert(ids=ids, documents=chunks, metadatas=metadatas)

    # Mark as ingested
    if file_id:
        _mark_ingested(file_id)

    print(f"  ✅ Ingested {len(chunks)} chunks from: {filename}")
    return len(chunks)


def query(question: str, stream: str, year: int, top_k: int = 5) -> list[dict]:
    """Query ChromaDB for relevant chunks, filtered by stream + year."""
    where_filter = {
        "$and": [
            {"stream": {"$eq": stream.upper()}},
            {"year": {"$eq": year}},
        ]
    }

    try:
        results = collection.query(
            query_texts=[question],
            n_results=top_k,
            where=where_filter,
        )
    except Exception as e:
        print(f"ChromaDB query error: {e}")
        return []

    # Format results
    documents = []
    if results and results["documents"] and results["documents"][0]:
        for i, doc in enumerate(results["documents"][0]):
            meta = results["metadatas"][0][i] if results["metadatas"] else {}
            documents.append({
                "content": doc,
                "filename": meta.get("filename", "Unknown"),
                "type": meta.get("type", "notes"),
                "drive_link": meta.get("drive_link", ""),
                "distance": results["distances"][0][i] if results["distances"] else 0,
            })

    return documents


def get_stats() -> dict:
    """Get stats about the vector store."""
    count = collection.count()
    return {"total_chunks": count}
